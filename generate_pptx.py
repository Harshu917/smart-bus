import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_simplified_sih_presentation(output_path="Smart_India_Hackathon_2025_Idea_Submission.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors
    NAVY_BLUE = RGBColor(31, 56, 100)      # #1F3864 Primary Accent
    FOOTER_BLUE = RGBColor(11, 94, 187)    # #0B5EBB Secondary Accent
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(25, 25, 25)           # Body text
    LIGHT_BG = RGBColor(248, 250, 253)
    BORDER_BLUE = RGBColor(210, 225, 245)

    def add_common_decorations(slide, slide_num, total_slides=6, show_team_badge=True):
        # 1. Top-Left Team Name Badge
        if show_team_badge:
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(0.4), Inches(2.2), Inches(0.45)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = LIGHT_BG
            badge.line.color.rgb = BORDER_BLUE
            badge.line.width = Pt(1.5)
            tf = badge.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Your Team Name"
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = NAVY_BLUE

        # 2. Top-Right SIH 2025 Logo Placeholder
        logo_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.7), Inches(0.4), Inches(2.833), Inches(0.5)
        )
        logo_box.fill.solid()
        logo_box.fill.fore_color.rgb = LIGHT_BG
        logo_box.line.color.rgb = BORDER_BLUE
        logo_box.line.width = Pt(1.5)
        tf_logo = logo_box.text_frame
        p_logo = tf_logo.paragraphs[0]
        p_logo.text = "💡 SMART INDIA HACKATHON 2025"
        p_logo.alignment = PP_ALIGN.CENTER
        p_logo.font.name = "Georgia"
        p_logo.font.size = Pt(10.5)
        p_logo.font.bold = True
        p_logo.font.color.rgb = NAVY_BLUE

        # 3. Bottom Footer Bar
        footer_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.95), Inches(13.333), Inches(0.55)
        )
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = FOOTER_BLUE
        footer_bar.line.fill.background()

        # Left Footer Text
        txBox_fl = slide.shapes.add_textbox(Inches(0.8), Inches(6.97), Inches(6.0), Inches(0.45))
        p_fl = txBox_fl.text_frame.paragraphs[0]
        p_fl.text = "@SIH Idea submission - Template"
        p_fl.font.name = "Arial"
        p_fl.font.size = Pt(11)
        p_fl.font.color.rgb = WHITE

        # Right Footer Slide Counter
        txBox_fr = slide.shapes.add_textbox(Inches(10.5), Inches(6.97), Inches(2.0), Inches(0.45))
        p_fr = txBox_fr.text_frame.paragraphs[0]
        p_fr.alignment = PP_ALIGN.RIGHT
        p_fr.text = f"{slide_num} / {total_slides}"
        p_fr.font.name = "Arial"
        p_fr.font.size = Pt(11)
        p_fr.font.bold = True
        p_fr.font.color.rgb = WHITE

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    add_common_decorations(s1, 1, show_team_badge=False)

    title_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(1.3))
    tf1 = title_box.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "SMART INDIA HACKATHON 2025"
    p1.font.name = "Georgia"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = NAVY_BLUE

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "TITLE PAGE"
    p1_sub.font.name = "Georgia"
    p1_sub.font.size = Pt(20)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = FOOTER_BLUE
    p1_sub.space_before = Pt(4)

    card1 = s1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.6), Inches(11.733), Inches(4.0)
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = LIGHT_BG
    card1.line.color.rgb = BORDER_BLUE
    card1.line.width = Pt(1.5)

    tf_c1 = card1.text_frame
    tf_c1.margin_left = Inches(0.4)
    tf_c1.margin_top = Inches(0.3)

    items_s1 = [
        ("Problem Statement ID", "SIH1608  (or Enter Official PS ID)"),
        ("Problem Statement Title", "Real-Time Smart Bus Tracking & Fleet Management System"),
        ("Theme", "Smart Automation / Transportation & Logistics / Smart Vehicles"),
        ("PS Category", "Software / Hardware"),
        ("Team ID", "SIH2025-TM-XXXXX  (Enter Registered Team ID)"),
        ("Team Name", "Team TransitPulse  (Registered on Portal)")
    ]

    for i, (label, val) in enumerate(items_s1):
        p = tf_c1.paragraphs[0] if i == 0 else tf_c1.add_paragraph()
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.space_before = Pt(10) if i > 0 else Pt(0)

        r_dot = p.add_run()
        r_dot.text = "• "
        r_dot.font.bold = True
        r_dot.font.color.rgb = FOOTER_BLUE

        r_lbl = p.add_run()
        r_lbl.text = f"{label} – "
        r_lbl.font.bold = True
        r_lbl.font.color.rgb = NAVY_BLUE

        r_val = p.add_run()
        r_val.text = val
        r_val.font.color.rgb = BLACK

    # =========================================================================
    # SLIDE 2: IDEA TITLE
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_common_decorations(s2, 2)

    h2 = s2.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(1.1))
    tf2 = h2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "IDEA TITLE"
    p2.font.name = "Georgia"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = NAVY_BLUE

    p2_sub = tf2.add_paragraph()
    p2_sub.text = "Proposed Solution (Describe your Idea/Solution/Prototype)"
    p2_sub.font.name = "Georgia"
    p2_sub.font.size = Pt(18)
    p2_sub.font.bold = True
    p2_sub.font.underline = True
    p2_sub.font.color.rgb = NAVY_BLUE
    p2_sub.space_before = Pt(3)

    c2 = s2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.4), Inches(11.733), Inches(4.2)
    )
    c2.fill.solid()
    c2.fill.fore_color.rgb = LIGHT_BG
    c2.line.color.rgb = BORDER_BLUE
    c2.line.width = Pt(1.5)

    tf_c2 = c2.text_frame
    tf_c2.margin_left = Inches(0.4)
    tf_c2.margin_top = Inches(0.3)

    s2_bullets = [
        ("Proposed Solution", "A real-time transit tracking platform with live GPS radar mapping, dynamic bus ETAs, and paperless digital QR ticketing."),
        ("How It Addresses the Problem", "Eliminates commuter waiting anxiety and overcrowding through instant bus location updates and live passenger occupancy counters."),
        ("Innovation & Uniqueness", "Sub-2-second WebSocket telemetry, hardware-agnostic design (works on smartphones or OBD-II/GPS devices), and offline-ready QR boarding passes.")
    ]

    for idx, (title, text) in enumerate(s2_bullets):
        p = tf_c2.paragraphs[0] if idx == 0 else tf_c2.add_paragraph()
        p.space_before = Pt(14) if idx > 0 else Pt(0)

        r_dot = p.add_run()
        r_dot.text = "• "
        r_dot.font.name = "Arial"
        r_dot.font.size = Pt(17)
        r_dot.font.bold = True
        r_dot.font.color.rgb = FOOTER_BLUE

        r_title = p.add_run()
        r_title.text = f"{title}: "
        r_title.font.name = "Arial"
        r_title.font.size = Pt(17)
        r_title.font.bold = True
        r_title.font.color.rgb = NAVY_BLUE

        r_txt = p.add_run()
        r_txt.text = text
        r_txt.font.name = "Arial"
        r_txt.font.size = Pt(15.5)
        r_txt.font.color.rgb = BLACK

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_common_decorations(s3, 3)

    h3 = s3.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(0.8))
    p3 = h3.text_frame.paragraphs[0]
    p3.text = "TECHNICAL APPROACH"
    p3.font.name = "Georgia"
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = NAVY_BLUE

    box_w = Inches(5.65)
    box_h = Inches(4.5)

    col1 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), box_w, box_h)
    col1.fill.solid()
    col1.fill.fore_color.rgb = LIGHT_BG
    col1.line.color.rgb = BORDER_BLUE
    col1.line.width = Pt(1.5)

    tf_c1 = col1.text_frame
    tf_c1.margin_left = Inches(0.3)
    tf_c1.margin_top = Inches(0.25)

    p_t1 = tf_c1.paragraphs[0]
    p_t1.text = "💻 Technologies Used"
    p_t1.font.name = "Georgia"
    p_t1.font.size = Pt(17)
    p_t1.font.bold = True
    p_t1.font.color.rgb = NAVY_BLUE

    techs = [
        ("Frontend", "React 19 + Tailwind CSS + Lucide Icons"),
        ("Live Map", "Leaflet GIS Radar + OpenStreetMap"),
        ("Backend", "Node.js + Express REST APIs"),
        ("Real-Time", "Socket.io (2-second WebSocket stream)"),
        ("Database", "In-Memory Data Store / PostgreSQL"),
        ("Hardware/IoT", "GPS Receiver (NavIC/GPS) + 4G LTE Modem")
    ]

    for label, val in techs:
        p = tf_c1.add_paragraph()
        p.space_before = Pt(8)
        r_b = p.add_run()
        r_b.text = "• "
        r_b.font.bold = True
        r_b.font.color.rgb = FOOTER_BLUE
        r_l = p.add_run()
        r_l.text = f"{label}: "
        r_l.font.bold = True
        r_l.font.size = Pt(14)
        r_l.font.color.rgb = NAVY_BLUE
        r_v = p.add_run()
        r_v.text = val
        r_v.font.size = Pt(13.5)
        r_v.font.color.rgb = BLACK

    col2 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.88), Inches(2.1), box_w, box_h)
    col2.fill.solid()
    col2.fill.fore_color.rgb = LIGHT_BG
    col2.line.color.rgb = BORDER_BLUE
    col2.line.width = Pt(1.5)

    tf_c2 = col2.text_frame
    tf_c2.margin_left = Inches(0.3)
    tf_c2.margin_top = Inches(0.25)

    p_t2 = tf_c2.paragraphs[0]
    p_t2.text = "⚙️ Process & Methodology"
    p_t2.font.name = "Georgia"
    p_t2.font.size = Pt(17)
    p_t2.font.bold = True
    p_t2.font.color.rgb = NAVY_BLUE

    steps = [
        ("1. Data Ingestion", "Bus beacon sends live GPS & speed every 2 seconds."),
        ("2. Real-Time Processing", "Server computes dynamic ETAs, delays, and stop geofences."),
        ("3. Instant Broadcast", "Coordinates pushed to Passenger and Driver screens via WebSockets."),
        ("4. Digital Ticketing", "Commuters reserve seats and generate scannable QR passes."),
        ("5. Fleet Analytics", "Admin monitors active buses, punctuality, and alerts.")
    ]

    for label, val in steps:
        p = tf_c2.add_paragraph()
        p.space_before = Pt(8)
        r_b = p.add_run()
        r_b.text = "✔ "
        r_b.font.bold = True
        r_b.font.color.rgb = FOOTER_BLUE
        r_l = p.add_run()
        r_l.text = f"{label}: "
        r_l.font.bold = True
        r_l.font.size = Pt(14)
        r_l.font.color.rgb = NAVY_BLUE
        r_v = p.add_run()
        r_v.text = val
        r_v.font.size = Pt(13.5)
        r_v.font.color.rgb = BLACK

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_common_decorations(s4, 4)

    h4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(0.8))
    p4 = h4.text_frame.paragraphs[0]
    p4.text = "FEASIBILITY AND VIABILITY"
    p4.font.name = "Georgia"
    p4.font.size = Pt(28)
    p4.font.bold = True
    p4.font.color.rgb = NAVY_BLUE

    c_w = Inches(3.7)
    c_h = Inches(4.5)
    gap = Inches(0.316)

    s4_data = [
        ("Feasibility Analysis", [
            ("Technical", "Built on lightweight open web standards; zero complex setup."),
            ("Operational", "Runs directly on existing driver smartphones or OBD-II units."),
            ("Economic", "Low cloud hosting cost; scales easily with bus fleet growth.")
        ], FOOTER_BLUE),
        ("Challenges & Risks", [
            ("Signal Drops", "Loss of GPS/cellular signal in tunnels and remote zones."),
            ("Peak Concurrency", "High traffic during morning and evening rush hours."),
            ("Driver Adoption", "Need simple interface with minimal driver distraction.")
        ], RGBColor(190, 40, 40)),
        ("Mitigation Strategies", [
            ("Dead Reckoning", "Client predicts location smoothly during temporary signal loss."),
            ("Socket Clustering", "Redis pub/sub handles 50,000+ simultaneous users."),
            ("1-Tap UI", "Large touch buttons, voice cues, and automated stop check-ins.")
        ], RGBColor(30, 140, 60))
    ]

    for idx, (title, items, col_border) in enumerate(s4_data):
        card = s4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8) + idx * (c_w + gap), Inches(2.1), c_w, c_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = col_border
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.name = "Georgia"
        pt.font.size = Pt(17)
        pt.font.bold = True
        pt.font.color.rgb = NAVY_BLUE

        for sub, desc in items:
            p = tf.add_paragraph()
            p.space_before = Pt(12)
            rb = p.add_run()
            rb.text = "• "
            rb.font.bold = True
            rb.font.color.rgb = col_border
            rl = p.add_run()
            rl.text = f"{sub}: "
            rl.font.bold = True
            rl.font.size = Pt(13.5)
            rl.font.color.rgb = NAVY_BLUE
            rv = p.add_run()
            rv.text = desc
            rv.font.size = Pt(13)
            rv.font.color.rgb = BLACK

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_common_decorations(s5, 5)

    h5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(0.8))
    p5 = h5.text_frame.paragraphs[0]
    p5.text = "IMPACT AND BENEFITS"
    p5.font.name = "Georgia"
    p5.font.size = Pt(28)
    p5.font.bold = True
    p5.font.color.rgb = NAVY_BLUE

    c1_5 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), box_w, box_h)
    c1_5.fill.solid()
    c1_5.fill.fore_color.rgb = LIGHT_BG
    c1_5.line.color.rgb = BORDER_BLUE
    c1_5.line.width = Pt(1.5)

    tf_c1_5 = c1_5.text_frame
    tf_c1_5.margin_left = Inches(0.3)
    tf_c1_5.margin_top = Inches(0.25)

    p_t1 = tf_c1_5.paragraphs[0]
    p_t1.text = "🎯 Impact on Target Audience"
    p_t1.font.name = "Georgia"
    p_t1.font.size = Pt(17)
    p_t1.font.bold = True
    p_t1.font.color.rgb = NAVY_BLUE

    audience = [
        ("Daily Commuters", "40% reduction in bus stop wait time; transparent real-time bus arrivals."),
        ("Transport Authorities", "Complete fleet visibility, automated route auditing, and delay tracking."),
        ("Bus Drivers", "Digital HUD console with stop checklists and 1-tap SOS emergency alert."),
        ("Women & Seniors", "Safer journeys with live trip sharing and contactless QR boarding.")
    ]

    for label, val in audience:
        p = tf_c1_5.add_paragraph()
        p.space_before = Pt(10)
        rb = p.add_run()
        rb.text = "★ "
        rb.font.bold = True
        rb.font.color.rgb = FOOTER_BLUE
        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.bold = True
        rl.font.size = Pt(14)
        rl.font.color.rgb = NAVY_BLUE
        rv = p.add_run()
        rv.text = val
        rv.font.size = Pt(13.5)
        rv.font.color.rgb = BLACK

    c2_5 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.88), Inches(2.1), box_w, box_h)
    c2_5.fill.solid()
    c2_5.fill.fore_color.rgb = LIGHT_BG
    c2_5.line.color.rgb = BORDER_BLUE
    c2_5.line.width = Pt(1.5)

    tf_c2_5 = c2_5.text_frame
    tf_c2_5.margin_left = Inches(0.3)
    tf_c2_5.margin_top = Inches(0.25)

    p_t2 = tf_c2_5.paragraphs[0]
    p_t2.text = "🌟 Solution Benefits"
    p_t2.font.name = "Georgia"
    p_t2.font.size = Pt(17)
    p_t2.font.bold = True
    p_t2.font.color.rgb = NAVY_BLUE

    benefits = [
        ("Social Impact", "Accessible, predictable, and dignified public transit for all citizens."),
        ("Economic Gains", "Reduces fleet fuel wastage from idling; boosts ticket revenue collection."),
        ("Environmental", "Encourages public bus adoption, reducing traffic congestion and urban emissions."),
        ("Smart Governance", "Provides data analytics for municipal transport planning and route optimization.")
    ]

    for label, val in benefits:
        p = tf_c2_5.add_paragraph()
        p.space_before = Pt(10)
        rb = p.add_run()
        rb.text = "✔ "
        rb.font.bold = True
        rb.font.color.rgb = FOOTER_BLUE
        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.bold = True
        rl.font.size = Pt(14)
        rl.font.color.rgb = NAVY_BLUE
        rv = p.add_run()
        rv.text = val
        rv.font.size = Pt(13.5)
        rv.font.color.rgb = BLACK

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_common_decorations(s6, 6)

    h6 = s6.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(0.8))
    p6 = h6.text_frame.paragraphs[0]
    p6.text = "RESEARCH AND REFERENCES"
    p6.font.name = "Georgia"
    p6.font.size = Pt(28)
    p6.font.bold = True
    p6.font.color.rgb = NAVY_BLUE

    c6 = s6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), Inches(11.733), Inches(4.5)
    )
    c6.fill.solid()
    c6.fill.fore_color.rgb = LIGHT_BG
    c6.line.color.rgb = BORDER_BLUE
    c6.line.width = Pt(1.5)

    tf_c6 = c6.text_frame
    tf_c6.margin_left = Inches(0.4)
    tf_c6.margin_top = Inches(0.3)

    refs = [
        ("Government Frameworks", "Ministry of Housing and Urban Affairs (MoHUA) – Smart Cities Mission & National Urban Transport Policy (NUTP)."),
        ("Transit Standards", "General Transit Feed Specification (GTFS Realtime) & Open Geospatial Consortium (OGC) GIS Mapping Standards."),
        ("Academic Research", "IEEE Intelligent Transportation Systems – 'Low-Latency WebSocket Architectures for Public Transit Telematics'."),
        ("Working Prototype", "Live Smart Bus prototype validated with 2-second telemetry broadcasts, radar map HUD, and QR ticketing.")
    ]

    for idx, (label, val) in enumerate(refs):
        p = tf_c6.paragraphs[0] if idx == 0 else tf_c6.add_paragraph()
        p.space_before = Pt(14) if idx > 0 else Pt(0)

        rb = p.add_run()
        rb.text = "• "
        rb.font.bold = True
        rb.font.size = Pt(16)
        rb.font.color.rgb = FOOTER_BLUE

        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.bold = True
        rl.font.size = Pt(16)
        rl.font.color.rgb = NAVY_BLUE

        rv = p.add_run()
        rv.text = val
        rv.font.size = Pt(15)
        rv.font.color.rgb = BLACK

    # Try saving to primary, and fallback to simple name if locked
    try:
        prs.save(output_path)
        print(f"Presentation saved to: {output_path}")
    except PermissionError:
        fallback = "Smart_India_Hackathon_2025_Idea_Submission_Simple.pptx"
        prs.save(fallback)
        print(f"Primary file was open in PowerPoint. Saved to: {fallback}")

if __name__ == "__main__":
    create_simplified_sih_presentation()
